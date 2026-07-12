import os
import logging
from abc import ABC, abstractmethod

from app.schemas.extraction import ExtractionResult

logger = logging.getLogger(__name__)


class BaseExtractor(ABC):

    @abstractmethod
    def extract(self, file_path: str) -> ExtractionResult:
        """Extract text from the file."""
        pass


class PDFExtractor(BaseExtractor):

    def extract(self, file_path: str) -> ExtractionResult:
        import fitz  # PyMuPDF

        text_content = []
        used_ocr = False

        try:
            with fitz.open(file_path) as doc:

                page_count = len(doc)

                # Standard extraction
                for page in doc:
                    text_content.append(page.get_text())

                combined_text = "\n".join(text_content).strip()

                # OCR fallback
                if not combined_text:
                    logger.info(f"[OCR START] {file_path}")

                    used_ocr = True
                    text_content = []

                    try:
                        for page in doc:
                            text_content.append(
                                page.get_text(
                                    "text",
                                    flags=fitz.TEXT_OCR
                                )
                            )

                        logger.info(f"[OCR SUCCESS] {file_path}")

                    except Exception as e:
                        logger.error(f"OCR failed: {e}")

                        return ExtractionResult(
                            success=False,
                            text="",
                            error=str(e),
                        )

            text = "\n".join(text_content)

            return ExtractionResult(
                success=True,
                text=text,
                page_count=page_count,
                word_count=len(text.split()),
                character_count=len(text),
                used_ocr=used_ocr,
                metadata={
                    "file_type": "pdf",
                    "filename": os.path.basename(file_path),
                    "pages": page_count,
                },
            )

        except Exception as e:

            return ExtractionResult(
                success=False,
                text="",
                error=str(e),
            )


class DocxExtractor(BaseExtractor):

    def extract(self, file_path: str) -> ExtractionResult:

        import docx

        try:
            doc = docx.Document(file_path)

            text_content = []

            for para in doc.paragraphs:
                if para.text.strip():
                    text_content.append(para.text)

            for table in doc.tables:
                for row in table.rows:
                    row_text = [
                        cell.text.strip()
                        for cell in row.cells
                        if cell.text.strip()
                    ]

                    if row_text:
                        text_content.append(" | ".join(row_text))

            text = "\n".join(text_content)

            return ExtractionResult(
                success=True,
                text=text,
                page_count=0,
                word_count=len(text.split()),
                character_count=len(text),
                used_ocr=False,
                metadata={
                    "file_type": "docx",
                    "filename": os.path.basename(file_path),
                },
            )

        except Exception as e:

            return ExtractionResult(
                success=False,
                text="",
                error=str(e),
            )


class PptxExtractor(BaseExtractor):

    def extract(self, file_path: str) -> ExtractionResult:

        import pptx

        try:

            prs = pptx.Presentation(file_path)

            text_content = []

            for i, slide in enumerate(prs.slides):

                slide_text = [f"--- Slide {i+1} ---"]

                for shape in slide.shapes:

                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())

                if len(slide_text) > 1:
                    text_content.append("\n".join(slide_text))

            text = "\n\n".join(text_content)

            return ExtractionResult(
                success=True,
                text=text,
                page_count=len(prs.slides),
                word_count=len(text.split()),
                character_count=len(text),
                used_ocr=False,
                metadata={
                    "file_type": "pptx",
                    "filename": os.path.basename(file_path),
                    "slides": len(prs.slides),
                },
            )

        except Exception as e:

            return ExtractionResult(
                success=False,
                text="",
                error=str(e),
            )


class XlsxExtractor(BaseExtractor):

    def extract(self, file_path: str) -> ExtractionResult:

        import openpyxl

        try:

            wb = openpyxl.load_workbook(
                file_path,
                read_only=True,
                data_only=True,
            )

            text_content = []

            for sheet_name in wb.sheetnames:

                sheet = wb[sheet_name]

                sheet_text = [f"--- Sheet: {sheet_name} ---"]

                for row in sheet.iter_rows(values_only=True):

                    row_vals = [
                        str(val).strip()
                        for val in row
                        if val is not None
                    ]

                    if row_vals:
                        sheet_text.append(" | ".join(row_vals))

                if len(sheet_text) > 1:
                    text_content.append("\n".join(sheet_text))

            text = "\n\n".join(text_content)

            return ExtractionResult(
                success=True,
                text=text,
                page_count=len(wb.sheetnames),
                word_count=len(text.split()),
                character_count=len(text),
                used_ocr=False,
                metadata={
                    "file_type": "xlsx",
                    "filename": os.path.basename(file_path),
                    "sheets": len(wb.sheetnames),
                },
            )

        except Exception as e:

            return ExtractionResult(
                success=False,
                text="",
                error=str(e),
            )


class TextExtractor(BaseExtractor):

    def extract(self, file_path: str) -> ExtractionResult:

        try:

            with open(
                file_path,
                "r",
                encoding="utf-8",
                errors="ignore",
            ) as f:

                text = f.read()

            return ExtractionResult(
                success=True,
                text=text,
                page_count=0,
                word_count=len(text.split()),
                character_count=len(text),
                used_ocr=False,
                metadata={
                    "file_type": "txt",
                    "filename": os.path.basename(file_path),
                },
            )

        except Exception as e:

            return ExtractionResult(
                success=False,
                text="",
                error=str(e),
            )

class ImageExtractor(BaseExtractor):
    def extract(self, file_path: str) -> ExtractionResult:
        try:
            from app.services.drawing_ocr import DrawingOCR
            text = DrawingOCR.extract_text([file_path])
            return ExtractionResult(
                success=True,
                text=text,
                page_count=1,
                word_count=len(text.split()),
                character_count=len(text),
                used_ocr=True,
                metadata={
                    "file_type": "image",
                    "filename": os.path.basename(file_path),
                },
            )
        except Exception as e:
            return ExtractionResult(
                success=False,
                text="",
                error=str(e),
            )

class DocumentExtractor:

    _extractors = {
        ".pdf": PDFExtractor(),
        ".docx": DocxExtractor(),
        ".pptx": PptxExtractor(),
        ".xlsx": XlsxExtractor(),
        ".txt": TextExtractor(),
        ".png": ImageExtractor(),
        ".jpg": ImageExtractor(),
        ".jpeg": ImageExtractor(),
    }

    @classmethod
    def extract(cls, file_path: str) -> ExtractionResult:

        _, ext = os.path.splitext(file_path)

        ext = ext.lower()

        extractor = cls._extractors.get(ext)

        if extractor is None:

            return ExtractionResult(
                success=False,
                text="",
                error=f"Unsupported file type: {ext}",
            )

        return extractor.extract(file_path)